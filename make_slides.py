"""Generate the presentation deck (.pptx) for the meta-labeling pipeline.

    python make_slides.py   ->   reports/Meta_Labeling_Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
GREY = RGBColor(0x55, 0x5B, 0x66)
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DATE = "June 18, 2026"
FOOT = "Multi-Asset Meta-Labeling  ·  " + DATE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_n = [0]


def _tb(slide, left, top, width, height, anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def _rect(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _header(slide, title):
    _rect(slide, 0.0, 0.0, 0.22, 7.5, NAVY)            # left spine
    t = _tb(slide, 0.6, 0.35, 12.4, 0.9)
    p = t.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY
    _rect(slide, 0.62, 1.18, 4.2, 0.045, ACCENT)       # accent underline


def _footer(slide):
    _n[0] += 1
    f = _tb(slide, 0.6, 7.02, 12.2, 0.35)
    p = f.paragraphs[0]; p.text = FOOT
    p.font.size = Pt(9); p.font.color.rgb = GREY
    pg = _tb(slide, 12.4, 7.02, 0.7, 0.35)
    p = pg.paragraphs[0]; p.text = str(_n[0]); p.font.size = Pt(9); p.font.color.rgb = GREY


def title_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0.0, 0.0, 13.333, 0.35, NAVY)
    _rect(s, 0.0, 7.15, 13.333, 0.35, ACCENT)
    tf = _tb(s, 0.9, 2.5, 11.5, 2.4)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(18); p2.font.color.rgb = GREY; p2.space_before = Pt(10)
    p3 = tf.add_paragraph(); p3.text = DATE
    p3.font.size = Pt(16); p3.font.bold = True; p3.font.color.rgb = ACCENT; p3.space_before = Pt(6)
    return s


def content(title, bullets, takeaway=None):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    body = _tb(s, 0.7, 1.45, 12.0, 4.7)
    for i, (txt, lvl) in enumerate(bullets):
        par = body.paragraphs[0] if i == 0 else body.add_paragraph()
        par.text = ("■  " if lvl == 0 else "–  ") + txt
        par.level = lvl
        par.font.size = Pt(17 if lvl == 0 else 14)
        par.font.bold = lvl == 0
        par.font.color.rgb = NAVY if lvl == 0 else GREY
        par.space_after = Pt(5)
    if takeaway:
        _rect(s, 0.6, 6.25, 12.13, 0.62, LIGHT)
        tk = _tb(s, 0.8, 6.3, 11.8, 0.55, MSO_ANCHOR.MIDDLE)
        p = tk.paragraphs[0]; p.text = "Takeaway:  " + takeaway
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = ACCENT
    _footer(s)
    return s


def table_slide(title, headers, rows, note=None, highlight_row=None, note_label="What it means"):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    nr, nc = len(rows) + 1, len(headers)
    rh = 0.34
    gt = s.shapes.add_table(nr, nc, Inches(0.7), Inches(1.5), Inches(11.9), Inches(rh * nr)).table
    for c, h in enumerate(headers):
        cell = gt.cell(0, c); cell.text = h
        pr = cell.text_frame.paragraphs[0]
        pr.font.size = Pt(13); pr.font.bold = True; pr.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        hot = highlight_row is not None and r - 1 == highlight_row
        for c, val in enumerate(row):
            cell = gt.cell(r, c); cell.text = str(val)
            pr = cell.text_frame.paragraphs[0]
            pr.font.size = Pt(13); pr.font.bold = hot
            pr.font.color.rgb = ACCENT if hot else NAVY
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if hot else WHITE
    if note:
        top = 1.5 + rh * nr + 0.25
        _rect(s, 0.6, top, 12.13, 1.0, LIGHT)
        tk = _tb(s, 0.8, top + 0.06, 11.8, 0.9, MSO_ANCHOR.MIDDLE)
        p = tk.paragraphs[0]; p.text = note_label + ":  " + note
        p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = ACCENT
    _footer(s)
    return s


# ---------------- slides ----------------

title_slide("Multi-Asset Meta-Labeling Research Pipeline",
            "A weekly, seven-sleeve global allocation framework  ·  research & educational use only")

content("Objective & Approach", [
    ("Allocate weekly across a diversified seven-sleeve global asset universe.", 0),
    ("Split the two decisions a single model usually blurs:", 0),
    ("M1 — direction: which assets to favour (simple, high recall)", 1),
    ("M2 — conviction: is the view likely to pay, and how much to size", 1),
    ("Meta-labeling design from systematic-PM literature (Lopez de Prado).", 0),
    ("Principle: STATIC factors live in M1; DYNAMIC factors live in M2.", 0),
], takeaway="Two simple problems beat one hard one — M1 proposes, M2 filters and sizes.")

content("Changes Since Last Review", [
    ("M1 simplified to a single static factor (momentum + trend); macro & risk moved out.", 0),
    ("ML layer kept to logistic regression only.", 0),
    ("Static factors → M1; dynamic factors (macro, volatility, regime) → M2.", 0),
    ("Macro / regime datasets now feed M2 (not used by earlier work).", 0),
    ("Added simple evaluations — mean difference, MAE, Brier, calibration MAPE.", 0),
    ("Out-of-sample confirmation — train/test split + 4-week embargo + walk-forward.", 0),
], takeaway="Every change follows one rule: static -> M1, dynamic -> M2, keep each layer simple.")

content("Architecture", [
    ("M1 — static, linear directional model (momentum + trend only).", 0),
    ("M2 — dynamic, regime-aware meta-label (logistic regression):", 0),
    ("sees momentum, trend, macro, volatility separately + regime features", 1),
    ("weights factors differently across regimes (factor-timing)", 1),
    ("Portfolio — benchmark-relative active weights, vol targeting, position caps, 2-layer costs.", 0),
    ("Discipline — no look-ahead, 4-week macro lag, 4-week train/test embargo.", 0),
], takeaway="M1 is simple on purpose; all dynamic / regime intelligence is isolated in M2.")

table_slide("M1 — Static Directional Model: worked 3-asset example",
            ["Asset", "12w mom (z)", "trend (z)", "Technical", "Rank", "Tilt", "Weight"],
            [["S&P 500", "+1.08", "+0.99", "+1.03", "best", "+10%", "43.3%"],
             ["Gold", "-0.90", "+0.02", "-0.44", "mid", "0%", "33.3%"],
             ["7-10Y Treasury", "-0.18", "-1.01", "-0.59", "worst", "-10%", "23.3%"]],
            note="Score each factor, z-score across assets, rank, then tilt around the 1/7 benchmark. "
                 "An underweight is an implicit short — no outright shorting needed.",
            highlight_row=0)

content("M2 — Dynamic Meta-Label", [
    ("Target — did M1's benchmark-relative bet pay over the next 4 weeks?", 0),
    ("Model — logistic regression (simple by design), refit on a rolling ~12-month window.", 0),
    ("Features — momentum, trend, macro, volatility (each separate) + regime (VIX, curve, credit, growth, inflation).", 0),
    ("Output — P(success), converted into a sizing multiplier on the M1 tilt.", 0),
    ("Evaluated with multiple methods — F1, AUC-ROC, AUC-PR, calibration — not just accuracy.", 0),
], takeaway="M2 is where 'which factor works in which regime' is meant to be learned.")

table_slide("Universe & Data",
            ["Sleeve", "Index", "Series used (free)", "History"],
            [["U.S. Equity", "S&P 500", "^GSPC (index)", "2000+"],
             ["Developed Intl", "MSCI EAFE", "EFA (ETF proxy)", "2001+"],
             ["Emerging Mkts", "MSCI EM", "EEM (ETF proxy)", "2003+"],
             ["U.S. Treasury 7-10Y", "S&P UST 7-10Y", "IEF (ETF proxy)", "2002+"],
             ["High Yield", "ICE BofA US HY", "HYG (ETF proxy)", "2007+"],
             ["Gold", "Gold spot", "GC=F (futures)", "2000+"],
             ["U.S. REITs", "Nasdaq US REIT", "FRED index", "2011+"]],
            note="Free true-index history is short (HY index only 2023, MSCI EAFE/EM only 2012); we use the "
                 "longest free series per sleeve. True long-history index needs Bloomberg.")

table_slide("Results — Strategy   (full sample 2000-2026; OOS from 2021)",
            ["Strategy", "Sharpe", "Max DD", "Sharpe (OOS)", "Info Ratio (OOS)"],
            [["Equal-Weight", "0.62", "-30%", "0.76", "—"],
             ["Moderate Growth", "0.56", "-34%", "0.71", "-0.32"],
             ["Institutional", "0.63", "-29%", "0.74", "-0.64"],
             ["M1-only", "0.65", "-22%", "0.81", "+0.18"],
             ["M1 + M2", "0.61", "-25%", "0.78", "-0.00"]],
            note="M1 beats all baselines on risk-adjusted return and cuts drawdown the most "
                 "(-22% vs -30%). Adding M2 makes every metric slightly worse.",
            highlight_row=3)

table_slide("Results — Walk-Forward   (Sharpe by window)",
            ["Window", "Equal-Weight", "M1-only", "M1 + M2"],
            [["2014-2016", "0.30", "0.05", "0.17"],
             ["2016-2018", "0.59", "0.53", "0.49"],
             ["2018-2020", "0.48", "0.15", "0.15"],
             ["2021-now", "0.76", "0.81", "0.78"]],
            note="M1's edge over equal-weight is concentrated in 2021+, not uniform. "
                 "Full-sample numbers flatter M1 than the regime-by-regime view.",
            highlight_row=3)

content("Key Finding — M2 Does Not Add Value Yet", [
    ("M2's probabilities carry no information — realized success is flat across every predicted bucket:", 0),
    ("predicted 0.14 → realized 0.44   |   predicted 0.82 → realized 0.43", 1),
    ("AUC-ROC = 0.50 full / 0.46 out-of-sample — i.e. random, slightly worse OOS.", 0),
    ("Simple error metrics agree — Brier 0.31, calibration MAPE ≈ 0.49 (probabilities ~49% off).", 1),
    ("M2 underperforms M1-only in every walk-forward window.", 0),
    ("Held under both proxy and real macro data — robust, not a tuning artifact.", 0),
], takeaway="The meta-label as specified (per-asset, 4-week, benchmark-relative) extracts no conditional signal.")

content("Attribution — Factors & Costs", [
    ("M1 decomposed into its two sub-signals (full-sample Sharpe):", 0),
    ("momentum 0.61  ·  trend 0.65  ·  combined technical 0.65", 1),
    ("Interaction is positive (+0.015) — momentum and trend reinforce, not redundant.", 0),
    ("Cost layers (annualized): gross 6.48% → net of expense 6.39% → net of all 6.19%.", 0),
    ("transaction cost (~29 bps) is a larger drag than the expense ratio (~9 bps)", 1),
], takeaway="The directional value is real and complementary; costs are modest and turnover-driven.")

content("Limitations", [
    ("Research-grade data (free index / ETF-proxy series); true long-history index needs Bloomberg.", 0),
    ("Historical simulation only — no capacity, market impact, borrow, or live execution.", 0),
    ("M1's edge is regime-concentrated (recent), not uniform across history.", 0),
    ("M2 is not yet usable; its prediction target needs reformulation.", 0),
    ("Some diagnostics are full-sample; production needs broader walk-forward / purged CV.", 0),
], takeaway="Solid on each step — we report what works and what does not, without overselling.")

content("Questions for Discussion", [
    ("M2 reformulation — which direction?", 0),
    ("a regime gate (scale total active risk down in stress), or", 1),
    ("a factor-timer (momentum vs trend by regime), or a different target / horizon?", 1),
    ("Is the 4-week, benchmark-relative success label the right thing to predict?", 0),
    ("M1's edge is recent-regime concentrated — acceptable, or do we need cross-regime robustness?", 0),
    ("Data — prioritize Bloomberg true-index history (pre-2012)? Which sleeves first?", 0),
    ("Universe — stay at 7, or add IG credit + broad commodity (9)?", 0),
    ("Benchmark / tracking-error budget for the information-ratio framing?", 0),
], takeaway="We are solid on M1 and the framework; M2's design is the open research decision.")

out = Path(__file__).resolve().parent / "reports" / "Meta_Labeling_Deck.pptx"
prs.save(out)
print(f"saved {_n[0] + 1} slides -> {out}")
